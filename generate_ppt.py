# -*- coding: utf-8 -*-
"""
生成《海外智能调奶器产品方案规划》决策汇报 PPT
主轴：智能化为主线，硬件为载体
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ========== 配色方案（3 色 + 灰） ==========
COLOR_PRIMARY = RGBColor(0x1A, 0x3A, 0x5C)      # 深蓝主色
COLOR_ACCENT = RGBColor(0xE8, 0x7A, 0x2E)       # 橙色强调
COLOR_TEXT = RGBColor(0x33, 0x33, 0x33)         # 深灰文字
COLOR_LIGHT = RGBColor(0xF5, 0xF5, 0xF5)        # 浅灰背景
COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_GREEN = RGBColor(0x2E, 0x8B, 0x57)        # 绿色（正向）
COLOR_RED = RGBColor(0xC0, 0x39, 0x2B)          # 红色（负向）

FONT_NAME = "PingFang SC"

# ========== 初始化演示文稿 ==========
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

SW = prs.slide_width
SH = prs.slide_height

# ========== 辅助函数 ==========
def add_blank_slide():
    return prs.slides.add_slide(prs.slide_layouts[6])

def set_bg(slide, color=COLOR_WHITE):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, left, top, width, height, color, line_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(0.5)
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    return shape

def add_text(slide, left, top, width, height, text, size=14, color=COLOR_TEXT,
             bold=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font=FONT_NAME):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(4)
    tf.margin_right = Pt(4)
    tf.margin_top = Pt(2)
    tf.margin_bottom = Pt(2)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = font
    return tb

def add_multiline(slide, left, top, width, height, lines, size=13, color=COLOR_TEXT,
                  bullet=False, line_spacing=1.2):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(6)
    tf.margin_right = Pt(6)
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        prefix = "• " if bullet else ""
        run = p.add_run()
        run.text = prefix + line
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.name = FONT_NAME
    return tb

def add_title_bar(slide, title, subtitle=None):
    """页面顶部标题栏"""
    add_rect(slide, 0, 0, SW, Inches(0.9), COLOR_PRIMARY)
    add_rect(slide, 0, Inches(0.9), SW, Pt(3), COLOR_ACCENT)
    add_text(slide, Inches(0.5), Inches(0.15), Inches(11), Inches(0.6),
             title, size=24, color=COLOR_WHITE, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    if subtitle:
        add_text(slide, Inches(0.5), Inches(0.55), Inches(11), Inches(0.3),
                 subtitle, size=12, color=RGBColor(0xCC, 0xDD, 0xEE))

def add_page_footer(slide, page_num, total=13):
    add_text(slide, Inches(0.5), Inches(7.15), Inches(6), Inches(0.3),
             "海外智能调奶器产品方案规划  |  机密", size=9, color=RGBColor(0x99,0x99,0x99))
    add_text(slide, Inches(11.5), Inches(7.15), Inches(1.5), Inches(0.3),
             f"{page_num} / {total}", size=9, color=RGBColor(0x99,0x99,0x99), align=PP_ALIGN.RIGHT)


# ========== Slide 1: 封面 ==========
def slide_cover():
    s = add_blank_slide()
    set_bg(s, COLOR_PRIMARY)
    # 装饰条
    add_rect(s, 0, Inches(3.0), SW, Pt(4), COLOR_ACCENT)
    add_rect(s, 0, Inches(4.5), SW, Pt(2), RGBColor(0x3A,0x5A,0x7C))
    # 主标题
    add_text(s, Inches(0.8), Inches(1.8), Inches(11), Inches(1.0),
             "海外智能调奶器", size=44, color=COLOR_WHITE, bold=True)
    add_text(s, Inches(0.8), Inches(2.7), Inches(11), Inches(0.6),
             "产品方案规划", size=32, color=COLOR_ACCENT, bold=True)
    # 副标题
    add_text(s, Inches(0.8), Inches(3.3), Inches(11), Inches(0.5),
             "智能化为主线 · 硬件为载体 · 数据建护城河", size=18, color=RGBColor(0xDD,0xEE,0xFF))
    # 元信息
    add_text(s, Inches(0.8), Inches(5.2), Inches(11), Inches(0.4),
             "视角：现有海外品牌扩品类  |  主推市场：北美 + 欧洲大陆",
             size=13, color=RGBColor(0xAA,0xBB,0xCC))
    add_text(s, Inches(0.8), Inches(5.6), Inches(11), Inches(0.4),
             "编制时间：2026年6月  |  依据：市场调研报告 + 用户评论挖掘报告",
             size=13, color=RGBColor(0xAA,0xBB,0xCC))
    # 底部一句话
    add_rect(s, 0, Inches(6.5), SW, Inches(1.0), RGBColor(0x12,0x2A,0x44))
    add_text(s, Inches(0.8), Inches(6.6), Inches(11.5), Inches(0.8),
             "“Baby Brezza sells a formula dispenser with WiFi. We sell a feeding intelligence system that happens to make bottles.”",
             size=14, color=COLOR_ACCENT, anchor=MSO_ANCHOR.MIDDLE, bold=True)

# ========== Slide 2: 市场机会 ==========
def slide_market():
    s = add_blank_slide()
    set_bg(s)
    add_title_bar(s, "市场机会：智能调奶器的蓝海在 L2-L5", "全球市场增长稳健，但当前'智能'仅停留在 L1 远程遥控")
    # 左侧：市场规模卡片
    add_rect(s, Inches(0.5), Inches(1.2), Inches(5.8), Inches(2.6), COLOR_LIGHT)
    add_text(s, Inches(0.7), Inches(1.3), Inches(5.4), Inches(0.4),
             "全球市场规模", size=14, color=COLOR_PRIMARY, bold=True)
    add_multiline(s, Inches(0.7), Inches(1.8), Inches(5.4), Inches(2.0), [
        "2024 年：$4.85 - 12 亿（不同口径）",
        "2032-2033 年：$12.5 - 24 亿",
        "CAGR：8.9% - 11.2%",
        "北美 39% / 欧洲 29% / 亚太 20%",
    ], size=14, bullet=True, line_spacing=1.4)
    # 右侧：智能化层级卡片
    add_rect(s, Inches(6.7), Inches(1.2), Inches(6.1), Inches(5.5), COLOR_LIGHT)
    add_text(s, Inches(6.9), Inches(1.3), Inches(5.7), Inches(0.4),
             "智能化层级模型 vs 市场覆盖", size=14, color=COLOR_PRIMARY, bold=True)
    # 层级表格
    levels = [
        ("L1 远程遥控", "App 出奶、预设、基础提醒", "Baby Brezza / Burabi", "红海"),
        ("L2 自适应精准", "称重校准、扫码识别、偏差自检", "全市场空白", "V1 核心"),
        ("L3 数据洞察", "喂养记录、生长曲线、库存管理", "全市场空白", "V1-V2 护城河"),
        ("L4 AI 个性化", "喂养建议、阶段切换、过敏识别", "全市场空白", "V2-V3 差异化"),
        ("L5 生态联动", "Owlet/Nanit、远程医疗、订阅", "全市场空白", "V3 愿景"),
    ]
    y = Inches(1.85)
    for lv, desc, cov, tag in levels:
        color = COLOR_RED if "红海" in tag else COLOR_GREEN
        tag_color = COLOR_RED if "红海" in tag else COLOR_ACCENT
        add_rect(s, Inches(6.9), y, Inches(5.7), Inches(0.85), COLOR_WHITE, COLOR_PRIMARY)
        add_text(s, Inches(7.0), y + Inches(0.05), Inches(1.6), Inches(0.35),
                 lv, size=11, color=COLOR_PRIMARY, bold=True)
        add_text(s, Inches(7.0), y + Inches(0.4), Inches(3.5), Inches(0.4),
                 desc, size=10, color=COLOR_TEXT)
        add_text(s, Inches(10.6), y + Inches(0.05), Inches(1.9), Inches(0.35),
                 cov, size=9, color=color, align=PP_ALIGN.RIGHT)
        add_text(s, Inches(10.6), y + Inches(0.4), Inches(1.9), Inches(0.4),
                 tag, size=10, color=tag_color, bold=True, align=PP_ALIGN.RIGHT)
        y += Inches(0.95)
    # 左下：核心结论
    add_rect(s, Inches(0.5), Inches(4.0), Inches(5.8), Inches(2.7), COLOR_PRIMARY)
    add_text(s, Inches(0.7), Inches(4.1), Inches(5.4), Inches(0.4),
             "核心机会判断", size=14, color=COLOR_ACCENT, bold=True)
    add_multiline(s, Inches(0.7), Inches(4.55), Inches(5.4), Inches(2.1), [
        "L1 已成红海，Baby Brezza 占据",
        "L2-L5 全部空白，真正蓝海",
        "新进入者应锚定 L2 起步",
        "用数据与算法构建长期护城河",
    ], size=13, color=COLOR_WHITE, bullet=True, line_spacing=1.5)
    add_page_footer(s, 2)

# ========== Slide 3: 用户痛点（评论挖掘） ==========
def slide_painpoints():
    s = add_blank_slide()
    set_bg(s)
    add_title_bar(s, "用户痛点：龙头的'伪智能'是最大软肋", "评论挖掘显示 Baby Brezza WiFi 版配比偏差 -17~-19%，35 起 CPSC 投诉含婴儿住院")
    # 三列痛点卡片
    cards = [
        ("配比精度崩塌", COLOR_RED, [
            "WiFi 版 2oz 瓶偏差 -17.1%",
            "WiFi 版 8oz 瓶偏差 -19.2%",
            "比基础版（-9%）更差",
            "35 起 CPSC 投诉",
            "含婴儿营养不良住院",
        ], "Consumer Reports 2024"),
        ("品牌识别仍靠客服", COLOR_ACCENT, [
            "Similac Advance 官网找不到",
            "需致电客服查条码数字",
            "易与其他型号混淆",
            "BabyExo 需手动算配比",
            "扫码识别全市场空白",
        ], "CR 评测"),
        ("智能功能名不副实", COLOR_RED, [
            "App 仅远程出奶+5预设",
            "无喂养数据记录",
            "无生长曲线 / AI 建议",
            "Tommee 家长要滤芯提醒未实现",
            "单账户单设备，无保姆权限",
        ], "Mother&Baby / 官方手册"),
    ]
    x = Inches(0.5)
    for title, c, items, src in cards:
        add_rect(s, x, Inches(1.2), Inches(4.0), Inches(5.0), COLOR_LIGHT)
        add_rect(s, x, Inches(1.2), Inches(4.0), Inches(0.6), c)
        add_text(s, x + Inches(0.2), Inches(1.25), Inches(3.6), Inches(0.5),
                 title, size=15, color=COLOR_WHITE, bold=True, anchor=MSO_ANCHOR.MIDDLE)
        add_multiline(s, x + Inches(0.2), Inches(2.0), Inches(3.6), Inches(3.8),
                      items, size=12, bullet=True, line_spacing=1.5)
        add_text(s, x + Inches(0.2), Inches(5.8), Inches(3.6), Inches(0.3),
                 "来源：" + src, size=9, color=RGBColor(0x99,0x99,0x99))
        x += Inches(4.2)
    # 底部一句话
    add_rect(s, Inches(0.5), Inches(6.4), Inches(12.3), Inches(0.6), COLOR_PRIMARY)
    add_text(s, Inches(0.7), Inches(6.45), Inches(12), Inches(0.5),
             "结论：当前所有竞品的'智能'仅停留在 L1 远程遥控层。配比自适应、数据记录、AI 推荐、生态联动——全部空白。",
             size=13, color=COLOR_WHITE, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    add_page_footer(s, 3)

# ========== Slide 4: 产品组合定位 ==========
def slide_positioning():
    s = add_blank_slide()
    set_bg(s)
    add_title_bar(s, "产品组合：智能喂养系统主机 + 便携款衍生", "主推产品 A 锚定 L2-L3，衍生产品 B 复用云端能力")
    # 产品 A 大卡
    add_rect(s, Inches(0.5), Inches(1.2), Inches(7.5), Inches(5.5), COLOR_LIGHT)
    add_rect(s, Inches(0.5), Inches(1.2), Inches(7.5), Inches(0.7), COLOR_PRIMARY)
    add_text(s, Inches(0.7), Inches(1.25), Inches(7.1), Inches(0.6),
             "产品 A：智能喂养系统主机  |  $169-199", size=18, color=COLOR_WHITE, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    add_multiline(s, Inches(0.7), Inches(2.1), Inches(7.1), Inches(4.5), [
        "定位：L2 自适应精准 + L3 数据洞察",
        "市场：北美 + 欧洲大陆同步（双电压 100-240V）",
        "对标：Baby Brezza $230-250，价格优势 15-30% + 智能代差",
        "差异化：称重式校准 + 扫码识别 + 喂养数据 + 生长曲线",
        "认证：FDA + CE 同步启动",
        "保修：2 年全球（vs Baby Brezza 1 年仅美加）",
        "V1 即覆盖 P0+P1 智能功能，建立数据基座",
    ], size=13, bullet=True, line_spacing=1.6)
    # 产品 B 小卡
    add_rect(s, Inches(8.3), Inches(1.2), Inches(4.5), Inches(2.6), COLOR_LIGHT)
    add_rect(s, Inches(8.3), Inches(1.2), Inches(4.5), Inches(0.6), COLOR_ACCENT)
    add_text(s, Inches(8.5), Inches(1.25), Inches(4.1), Inches(0.5),
             "产品 B：智能便携款  |  $99-129", size=14, color=COLOR_WHITE, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    add_multiline(s, Inches(8.5), Inches(1.95), Inches(4.1), Inches(1.8), [
        "V2 推出，承接便携场景",
        "复用 A 的云端品牌库与 App",
        "硬件小型化，智能能力不缩水",
        "场景：旅行、祖辈家、第二居所",
    ], size=11, bullet=True, line_spacing=1.4)
    # 商业模式卡
    add_rect(s, Inches(8.3), Inches(4.0), Inches(4.5), Inches(2.7), COLOR_PRIMARY)
    add_text(s, Inches(8.5), Inches(4.1), Inches(4.1), Inches(0.4),
             "商业模式演进", size=14, color=COLOR_ACCENT, bold=True)
    add_multiline(s, Inches(8.5), Inches(4.55), Inches(4.1), Inches(2.1), [
        "V1：硬件销售（$169-199）",
        "V2：+ App Pro 订阅（$39/年）",
        "V3：+ 奶粉订阅（$39-59/月）",
        "数据锁定提升迁移成本",
    ], size=12, color=COLOR_WHITE, bullet=True, line_spacing=1.5)
    add_page_footer(s, 4)

# ========== Slide 5: 产品 A 智能功能规格 ==========
def slide_features():
    s = add_blank_slide()
    set_bg(s)
    add_title_bar(s, "产品 A 智能功能规格", "L2 自适应精准 + L3 数据洞察为 V1 核心，L4-L5 为迭代愿景")
    # 表格
    rows = [
        ("层级", "功能", "规格", "优先级"),
        ("L2", "称重式配比校准", "0.1g 精度，偏差 >±3% 自动告警停机", "P0"),
        ("L2", "奶粉扫码识别", "云端品牌库 V1 覆盖 Top 50 品牌", "P0"),
        ("L2", "温度自适应", "红外测温 + 出奶前二次测温", "P0"),
        ("L3", "喂养记录自动生成", "时间/容量/品牌/阶段自动记录", "P0"),
        ("L3", "智能库存管理", "余量推算 + 一键补货 + 订阅", "P0"),
        ("L3", "生长曲线联动", "WHO 曲线 + 喂养趋势关联", "P1"),
        ("L3", "多用户权限", "主账户 + 5 操作员 + 操作日志", "P1"),
        ("L4", "AI 喂养建议", "基于生长数据推送建议（V2）", "P2"),
        ("L5", "Owlet/Nanit 联动", "可穿戴设备数据同步（V3）", "P2"),
    ]
    left = Inches(0.5)
    top = Inches(1.2)
    col_w = [Inches(1.2), Inches(3.2), Inches(6.4), Inches(1.5)]
    row_h = Inches(0.55)
    # 表头
    x = left
    for i, h in enumerate(rows[0]):
        add_rect(s, x, top, col_w[i], row_h, COLOR_PRIMARY)
        add_text(s, x + Inches(0.1), top, col_w[i] - Inches(0.2), row_h,
                 h, size=12, color=COLOR_WHITE, bold=True, anchor=MSO_ANCHOR.MIDDLE,
                 align=PP_ALIGN.CENTER if i in (0,3) else PP_ALIGN.LEFT)
        x += col_w[i]
    # 数据行
    y = top + row_h
    for r, row in enumerate(rows[1:]):
        x = left
        bg = COLOR_WHITE if r % 2 == 0 else COLOR_LIGHT
        for i, cell in enumerate(row):
            add_rect(s, x, y, col_w[i], row_h, bg, RGBColor(0xDD,0xDD,0xDD))
            color = COLOR_TEXT
            bold = False
            if i == 0:
                color = COLOR_PRIMARY
                bold = True
            if i == 3:
                if "P0" in cell:
                    color = COLOR_RED
                    bold = True
                elif "P1" in cell:
                    color = COLOR_ACCENT
                    bold = True
                else:
                    color = RGBColor(0x99,0x99,0x99)
            add_text(s, x + Inches(0.1), y, col_w[i] - Inches(0.2), row_h,
                     cell, size=11, color=color, bold=bold, anchor=MSO_ANCHOR.MIDDLE,
                     align=PP_ALIGN.CENTER if i in (0,3) else PP_ALIGN.LEFT)
            x += col_w[i]
        y += row_h
    # 底部说明
    add_text(s, Inches(0.5), Inches(6.6), Inches(12), Inches(0.4),
             "硬件从简：称重传感器 / 红外测温 / 条码扫描 / WiFi+BLE / 3.5\" LCD / 100-240V——硬件规格服从智能能力",
             size=11, color=RGBColor(0x66,0x66,0x66))
    add_page_footer(s, 5)

# ========== Slide 6: 差异化卖点矩阵 ==========
def slide_diff():
    s = add_blank_slide()
    set_bg(s)
    add_title_bar(s, "差异化卖点：真智能 vs 伪智能", "对标 Baby Brezza WiFi 版，10 个维度全面代差")
    rows = [
        ("维度", "Baby Brezza WiFi", "本方案产品 A", "强度"),
        ("配比方式", "容积式（偏差 -17~-19%）", "称重式（±3% + 自检告警）", "★★★★★"),
        ("奶粉识别", "手动查号 / 致电客服", "扫码 + 图像识别 + 云端库", "★★★★"),
        ("精度验证", "无第三方报告", "第三方实验室 + App 导出", "★★★★"),
        ("喂养数据", "无", "自动记录 + 可视化", "★★★★"),
        ("生长曲线", "无", "WHO 曲线 + 趋势关联", "★★★"),
        ("库存管理", "低水位提醒", "余量推算 + 一键补货", "★★★"),
        ("多用户", "单账户单设备", "主账户 + 5 操作员", "★★★"),
        ("电压", "120V 仅美加", "100-240V 全球", "★★"),
        ("价格 / 保修", "$249.99 / 1 年美加", "$169-199 / 2 年全球", "★★★"),
    ]
    left = Inches(0.5)
    top = Inches(1.2)
    col_w = [Inches(2.0), Inches(4.0), Inches(4.8), Inches(1.5)]
    row_h = Inches(0.5)
    # 表头
    x = left
    for i, h in enumerate(rows[0]):
        add_rect(s, x, top, col_w[i], row_h, COLOR_PRIMARY)
        add_text(s, x + Inches(0.1), top, col_w[i] - Inches(0.2), row_h,
                 h, size=12, color=COLOR_WHITE, bold=True, anchor=MSO_ANCHOR.MIDDLE,
                 align=PP_ALIGN.CENTER if i in (3,) else PP_ALIGN.LEFT)
        x += col_w[i]
    y = top + row_h
    for r, row in enumerate(rows[1:]):
        x = left
        bg = COLOR_WHITE if r % 2 == 0 else COLOR_LIGHT
        for i, cell in enumerate(row):
            add_rect(s, x, y, col_w[i], row_h, bg, RGBColor(0xDD,0xDD,0xDD))
            color = COLOR_TEXT
            if i == 1:
                color = COLOR_RED
            elif i == 2:
                color = COLOR_GREEN
                pass
            elif i == 3:
                color = COLOR_ACCENT
            add_text(s, x + Inches(0.1), y, col_w[i] - Inches(0.2), row_h,
                     cell, size=11, color=color, anchor=MSO_ANCHOR.MIDDLE,
                     align=PP_ALIGN.CENTER if i == 3 else PP_ALIGN.LEFT)
            x += col_w[i]
        y += row_h
    # 底部一句话
    add_rect(s, Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.6), COLOR_ACCENT)
    add_text(s, Inches(0.7), Inches(6.55), Inches(12), Inches(0.5),
             "一句话差异化：Baby Brezza 卖带 WiFi 的调奶器；我们卖会调奶的智能喂养系统。",
             size=14, color=COLOR_WHITE, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    add_page_footer(s, 6)

# ========== Slide 7: 技术路线（数据架构） ==========
def slide_tech():
    s = add_blank_slide()
    set_bg(s)
    add_title_bar(s, "技术路线：端-边-云三层智能架构", "端侧精准 + 边侧聚合 + 云侧智能，离线降级保证可用")
    # 三层架构图
    # 云端
    add_rect(s, Inches(1.0), Inches(1.3), Inches(11.3), Inches(1.6), COLOR_PRIMARY)
    add_text(s, Inches(1.2), Inches(1.4), Inches(11), Inches(0.4),
             "云端 Cloud", size=16, color=COLOR_ACCENT, bold=True)
    add_multiline(s, Inches(1.2), Inches(1.85), Inches(11), Inches(1.0), [
        "奶粉品牌库（Top 50 → Long tail）  |  AI 喂养建议模型（V2）  |  WHO 生长曲线基准",
        "用户账户 / 多用户权限 / 操作日志  |  订阅订单系统（V3）  |  COPPA/GDPR-K 合规层",
    ], size=11, color=COLOR_WHITE, line_spacing=1.4)
    # 连接线
    add_text(s, Inches(6.0), Inches(2.95), Inches(1.3), Inches(0.3),
             "↓ HTTPS/TLS 1.3", size=10, color=COLOR_TEXT, align=PP_ALIGN.CENTER, bold=True)
    # 边端 App
    add_rect(s, Inches(2.5), Inches(3.3), Inches(8.3), Inches(1.4), COLOR_ACCENT)
    add_text(s, Inches(2.7), Inches(3.4), Inches(8), Inches(0.4),
             "边端 App（iOS / Android）", size=16, color=COLOR_WHITE, bold=True)
    add_multiline(s, Inches(2.7), Inches(3.8), Inches(8), Inches(0.9), [
        "喂养记录可视化 / 生长曲线 / 库存管理  |  偏差自检报告导出",
        "多用户权限管理 / 操作日志  |  离线缓存 + 联网同步",
    ], size=11, color=COLOR_WHITE, line_spacing=1.4)
    add_text(s, Inches(6.0), Inches(4.75), Inches(1.3), Inches(0.3),
             "↓ WiFi / BLE", size=10, color=COLOR_TEXT, align=PP_ALIGN.CENTER, bold=True)
    # 端端
    add_rect(s, Inches(1.0), Inches(5.1), Inches(11.3), Inches(1.6), COLOR_GREEN)
    add_text(s, Inches(1.2), Inches(5.2), Inches(11), Inches(0.4),
             "端端 Device（机器本地，离线可用）", size=16, color=COLOR_WHITE, bold=True)
    add_multiline(s, Inches(1.2), Inches(5.65), Inches(11), Inches(1.0), [
        "称重传感器（0.1g）+ 红外测温 + 条码扫描  |  本地出奶控制",
        "本地配比自检告警（偏差 >±3% 停机）  |  LCD 触控离线操作",
    ], size=11, color=COLOR_WHITE, line_spacing=1.4)
    # 关键决策
    add_text(s, Inches(0.5), Inches(6.85), Inches(12), Inches(0.3),
             "关键决策：称重式（非容积式）| 扫码+图像双模识别 | WiFi主+BLE辅 | 云端AI推理 | COPPA/GDPR-K 架构层合规",
             size=10, color=RGBColor(0x66,0x66,0x66))
    add_page_footer(s, 7)

# ========== Slide 8: 认证合规路线 ==========
def slide_compliance():
    s = add_blank_slide()
    set_bg(s)
    add_title_bar(s, "认证合规路线：扩品类借力既有体系", "复用既有 FDA/CE 认证主体与测试机构关系，压缩周期")
    # 时间轴
    items = [
        ("M1-M3", "FDA 510(k) 准备\n（若归类为医疗器械辅助）", "复用既有认证主体", COLOR_PRIMARY),
        ("M4-M6", "FCC + CE 同步启动\n欧盟技术文件", "复用既有测试机构", COLOR_ACCENT),
        ("M7-M9", "认证获批 + 试产\nCOPPA/GDPR-K 合规", "法务团队前置", COLOR_PRIMARY),
        ("M10-M12", "北美首发\nUKCA/RCM 准备（V2）", "既有流程复用", COLOR_GREEN),
    ]
    x = Inches(0.6)
    for phase, content, leverage, c in items:
        add_rect(s, x, Inches(1.5), Inches(2.9), Inches(3.5), COLOR_LIGHT)
        add_rect(s, x, Inches(1.5), Inches(2.9), Inches(0.6), c)
        add_text(s, x + Inches(0.15), Inches(1.55), Inches(2.6), Inches(0.5),
                 phase, size=14, color=COLOR_WHITE, bold=True, anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)
        add_text(s, x + Inches(0.15), Inches(2.3), Inches(2.6), Inches(1.6),
                 content, size=12, color=COLOR_TEXT, align=PP_ALIGN.CENTER)
        add_rect(s, x + Inches(0.15), Inches(4.0), Inches(2.6), Inches(0.9), COLOR_WHITE, c)
        add_text(s, x + Inches(0.15), Inches(4.05), Inches(2.6), Inches(0.85),
                 "借力点\n" + leverage, size=11, color=c, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        x += Inches(3.1)
    # 风险提示
    add_rect(s, Inches(0.6), Inches(5.3), Inches(12.1), Inches(1.5), COLOR_LIGHT)
    add_text(s, Inches(0.8), Inches(5.4), Inches(11.7), Inches(0.4),
             "关键风险与应对", size=14, color=COLOR_RED, bold=True)
    add_multiline(s, Inches(0.8), Inches(5.8), Inches(11.7), Inches(1.0), [
        "配比精度风险：称重式 + 第三方实验室验证 + 产品责任险 + CPSC 主动报备",
        "数据隐私：架构层 COPPA/GDPR-K 合规；数据最小化；欧盟数据不出欧",
        "AI 建议误判：V2 明确标注'参考性非医疗建议' + 儿科医生顾问委员会",
    ], size=11, bullet=True, line_spacing=1.4)
    add_page_footer(s, 8)

# ========== Slide 9: 迭代路线图 ==========
def slide_roadmap():
    s = add_blank_slide()
    set_bg(s)
    add_title_bar(s, "迭代路线图：V1 智能基座 → V2 AI 个性化 → V3 生态联动", "24-30 个月从 L2 走到 L5")
    # 三阶段
    phases = [
        ("V1 MVP", "0-12 个月", COLOR_PRIMARY, [
            "L2 全部 P0（称重+扫码+测温+自检）",
            "L3 全部 P0（喂养记录+库存管理）",
            "L3 P1（生长曲线+多用户+预警）",
            "北美首发，FDA+FCC+CE",
            "目标：配比±3% / 评分≥4.4",
        ]),
        ("V2", "+6 个月", COLOR_ACCENT, [
            "便携款 B 推出",
            "L4 AI 喂养建议引擎",
            "L4 过敏模式识别",
            "L2 图像识别升级",
            "欧洲大陆上架 / UKCA+RCM",
        ]),
        ("V3", "+12 个月", COLOR_GREEN, [
            "L5 Owlet/Nanit 联动",
            "L5 远程医疗数据导出（FHIR）",
            "L5 订阅制奶粉自动补货",
            "设备低价$129 + 奶粉订阅",
            "数据锁定护城河成型",
        ]),
    ]
    x = Inches(0.5)
    for name, period, c, items in phases:
        add_rect(s, x, Inches(1.2), Inches(4.1), Inches(5.3), COLOR_LIGHT)
        add_rect(s, x, Inches(1.2), Inches(4.1), Inches(0.8), c)
        add_text(s, x + Inches(0.2), Inches(1.25), Inches(3.7), Inches(0.4),
                 name, size=20, color=COLOR_WHITE, bold=True)
        add_text(s, x + Inches(0.2), Inches(1.6), Inches(3.7), Inches(0.35),
                 period, size=13, color=COLOR_WHITE)
        add_multiline(s, x + Inches(0.2), Inches(2.2), Inches(3.7), Inches(4.2),
                      items, size=12, bullet=True, line_spacing=1.7)
        x += Inches(4.2)
    # 底部里程碑
    add_text(s, Inches(0.5), Inches(6.7), Inches(12), Inches(0.4),
             "关键里程碑：M1-M3 配比方案验证 → M4-M6 App+云端库 → M7-M9 认证+试产 → M10-M12 北美首发",
             size=11, color=RGBColor(0x66,0x66,0x66))
    add_page_footer(s, 9)

# ========== Slide 10: 定价与渠道 ==========
def slide_pricing():
    s = add_blank_slide()
    set_bg(s)
    add_title_bar(s, "定价与渠道：硬件切入 + 数据锁定 + 订阅变现", "扩品类借力既有渠道/客服/KOL 资源")
    # 左：定价表
    add_rect(s, Inches(0.5), Inches(1.2), Inches(6.0), Inches(5.5), COLOR_LIGHT)
    add_text(s, Inches(0.7), Inches(1.3), Inches(5.6), Inches(0.4),
             "定价体系", size=16, color=COLOR_PRIMARY, bold=True)
    prices = [
        ("产品 A 主机", "$169-199", "对标 BB $230-250，毛利 35%+", COLOR_PRIMARY),
        ("产品 B 便携款", "$99-129", "对标 BB Mini $99，毛利 30%+", COLOR_ACCENT),
        ("App 基础版", "免费", "数据积累入口", RGBColor(0x66,0x66,0x66)),
        ("App Pro（V2）", "$39/年", "AI 建议，毛利 80%+", COLOR_GREEN),
        ("奶粉订阅（V3）", "$39-59/月", "与奶粉品牌合作，毛利 15-20%", COLOR_GREEN),
    ]
    y = Inches(1.8)
    for name, price, note, c in prices:
        add_rect(s, Inches(0.7), y, Inches(5.6), Inches(0.85), COLOR_WHITE, RGBColor(0xDD,0xDD,0xDD))
        add_rect(s, Inches(0.7), y, Inches(0.15), Inches(0.85), c)
        add_text(s, Inches(1.0), y + Inches(0.05), Inches(2.8), Inches(0.35),
                 name, size=12, color=COLOR_TEXT, bold=True)
        add_text(s, Inches(1.0), y + Inches(0.4), Inches(3.8), Inches(0.4),
                 note, size=10, color=RGBColor(0x66,0x66,0x66))
        add_text(s, Inches(4.0), y + Inches(0.05), Inches(2.2), Inches(0.75),
                 price, size=16, color=c, bold=True, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
        y += Inches(0.95)
    # 右：渠道策略
    add_rect(s, Inches(6.8), Inches(1.2), Inches(6.0), Inches(5.5), COLOR_PRIMARY)
    add_text(s, Inches(7.0), Inches(1.3), Inches(5.6), Inches(0.4),
             "渠道策略（扩品类借力）", size=16, color=COLOR_ACCENT, bold=True)
    channels = [
        ("V1 北美首发", "Amazon US + 独立站 DTC", "既有店铺评分 + DTC 流量"),
        ("V1-V2 欧洲大陆", "Amazon DE/FR/IT/ES + 本地化", "既有 EN/DE/FR 客服"),
        ("V2 线下拓展", "BuyBuyBaby / Boots / M&P", "既有线下母婴连锁关系"),
        ("全阶段营销", "KOL 资源库 + 邮件列表", "已购用户交叉销售"),
    ]
    y = Inches(1.85)
    for stage, ch, lev in channels:
        add_rect(s, Inches(7.0), y, Inches(5.6), Inches(1.05), RGBColor(0x12,0x2A,0x44))
        add_text(s, Inches(7.15), y + Inches(0.05), Inches(5.3), Inches(0.35),
                 stage, size=12, color=COLOR_ACCENT, bold=True)
        add_text(s, Inches(7.15), y + Inches(0.4), Inches(5.3), Inches(0.35),
                 "→ " + ch, size=11, color=COLOR_WHITE)
        add_text(s, Inches(7.15), y + Inches(0.7), Inches(5.3), Inches(0.3),
                 "借力：" + lev, size=10, color=RGBColor(0xAA,0xBB,0xCC))
        y += Inches(1.15)
    add_page_footer(s, 10)

# ========== Slide 11: 风险与应对 ==========
def slide_risks():
    s = add_blank_slide()
    set_bg(s)
    add_title_bar(s, "风险与应对", "智能化主轴下的五类核心风险")
    risks = [
        ("配比精度引发安全事件", "🔴🔴🔴🔴🔴",
         "称重式 + 第三方实验室验证 + 产品责任险 + CPSC 主动报备", COLOR_RED),
        ("数据隐私违规 COPPA/GDPR-K", "🔴🔴🔴🔴",
         "架构层合规；数据最小化；父母同意流程；欧盟数据不出欧", COLOR_RED),
        ("AI 建议误判", "🔴🔴🔴",
         "V2 标注'参考性非医疗建议' + 免责声明 + 儿科医生顾问委员会", COLOR_ACCENT),
        ("云端品牌库过时", "🔴🔴",
         "运营团队持续维护 + 用户反馈通道 + 未收录 24h 回传承诺", COLOR_ACCENT),
        ("Baby Brezza 降价防御", "🔴🔴",
         "以智能代差对抗价格战 + 数据锁定提升迁移成本", COLOR_ACCENT),
    ]
    y = Inches(1.3)
    for name, sev, response, c in risks:
        add_rect(s, Inches(0.5), y, Inches(12.3), Inches(1.0), COLOR_LIGHT)
        add_rect(s, Inches(0.5), y, Inches(0.2), Inches(1.0), c)
        add_text(s, Inches(0.85), y + Inches(0.1), Inches(4.5), Inches(0.4),
                 name, size=14, color=COLOR_TEXT, bold=True)
        add_text(s, Inches(0.85), y + Inches(0.5), Inches(4.5), Inches(0.4),
                 "严重度：" + sev, size=11, color=c)
        add_text(s, Inches(5.6), y + Inches(0.1), Inches(7.0), Inches(0.4),
                 "应对", size=11, color=COLOR_PRIMARY, bold=True)
        add_text(s, Inches(5.6), y + Inches(0.4), Inches(7.0), Inches(0.55),
                 response, size=11, color=COLOR_TEXT)
        y += Inches(1.1)
    add_page_footer(s, 11)

# ========== Slide 12: 扩品类借力清单 ==========
def slide_leverage():
    s = add_blank_slide()
    set_bg(s)
    add_title_bar(s, "扩品类借力：6 大基础设施复用", "不是从 0 建品牌，而是借用既有资源压缩 GTM 3-6 个月")
    levs = [
        ("渠道复用", "既有 Amazon 店铺 / 独立站 / 线下母婴连锁关系直接上架", "省去渠道冷启动 3-6 个月"),
        ("合规复用", "既有 FDA/CE 认证主体 + 测试机构关系 + 合规团队", "缩短 510(k)/CE 周期"),
        ("客服复用", "既有本地客服团队（EN/DE/FR）承接调奶器售后", "降低母婴退货率冲击"),
        ("品牌信任迁移", "以'母品牌×智能喂养'子系列推出", "借用既有信任度降低决策门槛"),
        ("数据复用", "既有用户画像 + KOL 资源库 + 邮件列表", "定向触达已购母婴用户交叉销售"),
        ("成本优势放大", "既有供应链/物流/海外仓关系", "使 $169-199 定价带保持 35%+ 毛利"),
    ]
    positions = [
        (Inches(0.5), Inches(1.3)),
        (Inches(4.7), Inches(1.3)),
        (Inches(8.9), Inches(1.3)),
        (Inches(0.5), Inches(3.6)),
        (Inches(4.7), Inches(3.6)),
        (Inches(8.9), Inches(3.6)),
    ]
    for (name, desc, benefit), (x, y) in zip(levs, positions):
        add_rect(s, x, y, Inches(4.0), Inches(2.1), COLOR_LIGHT)
        add_rect(s, x, y, Inches(4.0), Inches(0.55), COLOR_PRIMARY)
        add_text(s, x + Inches(0.15), y + Inches(0.05), Inches(3.7), Inches(0.45),
                 name, size=14, color=COLOR_WHITE, bold=True, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x + Inches(0.15), y + Inches(0.7), Inches(3.7), Inches(0.9),
                 desc, size=11, color=COLOR_TEXT, line_spacing=1.3)
        add_rect(s, x + Inches(0.15), y + Inches(1.65), Inches(3.7), Inches(0.35), COLOR_ACCENT)
        add_text(s, x + Inches(0.15), y + Inches(1.65), Inches(3.7), Inches(0.35),
                 "→ " + benefit, size=10, color=COLOR_WHITE, bold=True, anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)
    add_page_footer(s, 12)

# ========== Slide 13: 总结与下一步 ==========
def slide_summary():
    s = add_blank_slide()
    set_bg(s, COLOR_PRIMARY)
    # 标题
    add_text(s, Inches(0.8), Inches(0.6), Inches(11), Inches(0.6),
             "总结与下一步", size=32, color=COLOR_WHITE, bold=True)
    add_rect(s, Inches(0.8), Inches(1.2), Inches(4), Pt(3), COLOR_ACCENT)
    # 核心结论
    add_text(s, Inches(0.8), Inches(1.5), Inches(11), Inches(0.5),
             "方案核心结论", size=18, color=COLOR_ACCENT, bold=True)
    conclusions = [
        "市场蓝海确认：L2-L5 智能能力全市场空白，L1 已成红海",
        "产品定位：智能喂养系统主机 $169-199，锚定 L2-L3，硬件服从智能",
        "差异化武器：称重式校准 + 扫码识别 + 喂养数据 + 生长曲线",
        "迭代节奏：V1（12 月）→ V2（+6 月）→ V3（+12 月），24-30 月走到 L5",
        "扩品类借力：复用渠道/合规/客服/品牌/KOL/供应链 6 大基础设施",
    ]
    add_multiline(s, Inches(0.8), Inches(2.0), Inches(11.5), Inches(2.5),
                  conclusions, size=14, color=COLOR_WHITE, bullet=True, line_spacing=1.7)
    # 下一步
    add_rect(s, Inches(0.8), Inches(4.7), Inches(11.5), Inches(2.2), RGBColor(0x12,0x2A,0x44))
    add_text(s, Inches(1.0), Inches(4.8), Inches(11), Inches(0.4),
             "下一步行动（决策通过后 90 天）", size=16, color=COLOR_ACCENT, bold=True)
    nexts = [
        "D1-D30：组建智能产品团队（PM + 算法 + App + 硬件）；称重式配比方案 PoC",
        "D31-D60：第三方实验室精度测试；云端品牌库 Top 50 搭建启动",
        "D61-D90：FDA/FCC 认证启动；App 原型开发；KOL 种子用户招募",
    ]
    add_multiline(s, Inches(1.0), Inches(5.25), Inches(11), Inches(1.6),
                  nexts, size=12, color=COLOR_WHITE, bullet=True, line_spacing=1.5)
    # 决策请求
    add_rect(s, 0, Inches(7.0), SW, Inches(0.5), COLOR_ACCENT)
    add_text(s, Inches(0.8), Inches(7.05), Inches(11.5), Inches(0.4),
             "请决策：是否批准 V1 立项 + 预算 + 团队组建？",
             size=14, color=COLOR_WHITE, bold=True, anchor=MSO_ANCHOR.MIDDLE)


# ========== 生成所有 slide ==========
slide_cover()
slide_market()
slide_painpoints()
slide_positioning()
slide_features()
slide_diff()
slide_tech()
slide_compliance()
slide_roadmap()
slide_pricing()
slide_risks()
slide_leverage()
slide_summary()

# ========== 保存 ==========
output_path = "/Users/lute/WorkBuddy/2026-06-24-18-13-46/决策汇报.pptx"
prs.save(output_path)
print(f"PPT 已生成：{output_path}")
print(f"共 {len(prs.slides)} 张 slide")
